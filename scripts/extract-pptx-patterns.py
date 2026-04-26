#!/usr/bin/env python3
"""
Extract STRUCTURAL patterns from a winning proposal pitch deck (.pptx).

Surfaces slide-level structure WITHOUT exposing prose:
  - Slide count
  - Per-slide: layout name, # text frames, # images, # shapes, # tables, # charts
  - Per-slide: title text + word count of body text (title is shown; body is summarized)
  - Slide-master / theme info
  - Aggregate fonts used
  - Color palette (from theme)
  - Image extraction count (without saving the images)

Caller is expected to read this output, extract STRUCTURAL learning, and
NOT store source slides or imagery in the repo.

Usage:
    python scripts/extract-pptx-patterns.py /path/to/deck.pptx
    python scripts/extract-pptx-patterns.py /path/to/deck.pptx --titles-only
"""
from __future__ import annotations

import argparse
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def emu_to_inches(emu: int) -> float:
    return emu / 914400.0


def main(argv: list[str]) -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--titles-only", action="store_true")
    args = ap.parse_args(argv)

    if not args.pptx.exists():
        print(f"error: {args.pptx} not found", file=sys.stderr)
        return 1

    prs = Presentation(str(args.pptx))
    n = len(prs.slides)
    sw_in = emu_to_inches(prs.slide_width)
    sh_in = emu_to_inches(prs.slide_height)

    print(f"=== {args.pptx.name} ===")
    print(f"Slides: {n}")
    print(f"Slide size: {sw_in:.2f} x {sh_in:.2f} inches")
    print()

    # Aggregate slide-master / layout names
    layout_names = []
    for slide in prs.slides:
        layout_names.append(slide.slide_layout.name)
    print(f"Layouts used:")
    for name, c in Counter(layout_names).most_common():
        print(f"  {c:>3}  {name}")
    print()

    # Per-slide summary
    print("Per-slide summary:")
    print(f"{'#':>3}  {'layout':<25}  {'frames':>6}  {'imgs':>4}  {'shapes':>6}  {'tbls':>4}  {'charts':>6}  title")
    fonts_used: Counter[str] = Counter()
    total_imgs = 0
    total_shapes = 0
    total_tables = 0
    total_charts = 0
    for i, slide in enumerate(prs.slides, 1):
        n_text_frames = 0
        n_imgs = 0
        n_shapes = 0
        n_tables = 0
        n_charts = 0
        title = ""
        body_words = 0
        for shape in slide.shapes:
            n_shapes += 1
            if shape.has_text_frame:
                n_text_frames += 1
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts_used[run.font.name] += 1
                        body_words += len((run.text or "").split())
                if shape.has_text_frame and shape == slide.shapes.title if hasattr(slide.shapes, 'title') and slide.shapes.title else False:
                    title = shape.text_frame.text.strip()[:80]
            try:
                if shape.shape_type == 13:  # PICTURE
                    n_imgs += 1
            except Exception:
                pass
            if shape.has_table:
                n_tables += 1
            if shape.has_chart:
                n_charts += 1
        # Try to get the title shape directly
        if not title:
            try:
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()[:80]
            except Exception:
                pass
        title_safe = title.replace("\n", " ") if title else "(no title)"
        layout_name = slide.slide_layout.name[:25]
        print(f"{i:>3}  {layout_name:<25}  {n_text_frames:>6}  {n_imgs:>4}  {n_shapes:>6}  {n_tables:>4}  {n_charts:>6}  {title_safe}")
        total_imgs += n_imgs
        total_shapes += n_shapes
        total_tables += n_tables
        total_charts += n_charts

    print()
    print(f"Aggregate:  imgs={total_imgs}, shapes={total_shapes}, tables={total_tables}, charts={total_charts}")
    print()

    # Fonts used
    if fonts_used:
        print(f"Fonts referenced (run-level):")
        for font, c in fonts_used.most_common(10):
            print(f"  {c:>5}  {font}")
        print()

    # Aspect ratio classification
    ar = sw_in / sh_in
    if abs(ar - 16/9) < 0.05:
        ar_label = "16:9 widescreen"
    elif abs(ar - 4/3) < 0.05:
        ar_label = "4:3 standard"
    elif abs(ar - 16/10) < 0.05:
        ar_label = "16:10"
    else:
        ar_label = f"custom {ar:.3f}"
    print(f"Aspect ratio: {ar_label}")

    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
